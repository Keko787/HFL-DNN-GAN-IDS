"""Extract slide text, tables, notes and picture inventory from a .pptx."""
from __future__ import annotations

import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

if len(sys.argv) < 2:
    raise SystemExit("usage: python read_pptx.py <deck.pptx>  [> dump.txt]")
SRC = Path(sys.argv[1])


def walk(shapes, depth=0):
    out = []
    for sh in shapes:
        pad = "  " * depth
        try:
            if sh.shape_type is not None and str(sh.shape_type).startswith("GROUP"):
                out.append(f"{pad}[group]")
                out.extend(walk(sh.shapes, depth + 1))
                continue
        except Exception:
            pass
        if sh.has_table:
            t = sh.table
            out.append(f"{pad}[table {len(t.rows)}x{len(t.columns)}]")
            for r in t.rows:
                cells = [c.text.strip().replace("\n", " / ") for c in r.cells]
                out.append(pad + "  | " + " | ".join(cells))
            continue
        if sh.has_chart:
            ch = sh.chart
            out.append(f"{pad}[chart {ch.chart_type}]")
            try:
                cats = list(ch.plots[0].categories)
                out.append(f"{pad}  categories: {cats}")
                for s in ch.plots[0].series:
                    out.append(f"{pad}  series {s.name}: {list(s.values)}")
            except Exception:
                pass
            continue
        if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            w, h = Emu(sh.width).inches, Emu(sh.height).inches
            out.append(f"{pad}[picture {w:.1f}in x {h:.1f}in]")
            continue
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt:
                for i, line in enumerate(txt.split("\n")):
                    if line.strip():
                        out.append(f"{pad}{'· ' if i else ''}{line.strip()}")
    return out


def main() -> int:
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    prs = Presentation(str(SRC))
    print(f"FILE: {SRC.name}")
    print(f"SLIDES: {len(prs.slides)}   "
          f"CANVAS: {Emu(prs.slide_width).inches:.2f}in x "
          f"{Emu(prs.slide_height).inches:.2f}in")
    for i, slide in enumerate(prs.slides, 1):
        print(f"\n{'='*78}\n### SLIDE {i}\n{'='*78}")
        for line in walk(slide.shapes):
            print(line)
        if slide.has_notes_slide:
            note = slide.notes_slide.notes_text_frame.text.strip()
            if note:
                print(f"\n[NOTES] {note}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
